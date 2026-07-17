"""Geração de apresentação PPTX a partir de documentos."""

from __future__ import annotations

import json
import re
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from core.documentos import extrair_texto_arquivo
from core.openai_client import chat_completion, get_api_key
from core.utils import OUTPUTS_DIR, ensure_dirs

SYSTEM_PPTX = """Você monta outlines de slides executivos em português do Brasil.
Responda SOMENTE JSON válido no formato:
{
  "titulo": "...",
  "subtitulo": "...",
  "slides": [
    {"titulo": "...", "bullets": ["...", "..."]}
  ]
}
Regras: 6 a 10 slides; bullets curtos; não invente fatos fora das fontes."""


def _extrair_json(texto: str) -> dict[str, Any]:
    texto = texto.strip()
    if texto.startswith("```"):
        texto = re.sub(r"^```(?:json)?\s*", "", texto)
        texto = re.sub(r"\s*```$", "", texto)
    return json.loads(texto)


def _texto_fontes(caminhos: list[Path], *, max_chars: int = 60000) -> str:
    partes: list[str] = []
    total = 0
    for path in caminhos:
        try:
            doc = extrair_texto_arquivo(path.name, path.read_bytes())
        except Exception as exc:  # noqa: BLE001
            partes.append(f"--- {path.name} (erro: {exc}) ---")
            continue
        bloco = f"--- {path.name} ---\n{doc.texto}"
        if total + len(bloco) > max_chars:
            resto = max_chars - total
            if resto > 500:
                partes.append(bloco[:resto] + "\n…")
            break
        partes.append(bloco)
        total += len(bloco)
    return "\n\n".join(partes)


def gerar_outline_slides(caminhos: list[Path]) -> dict[str, Any]:
    if not get_api_key():
        raise RuntimeError("OPENAI_API_KEY não configurada.")
    fontes = _texto_fontes(caminhos)
    if not fontes.strip():
        raise ValueError("Fontes vazias.")
    resp = chat_completion(
        [
            {"role": "system", "content": SYSTEM_PPTX},
            {
                "role": "user",
                "content": (
                    "Com base nas fontes abaixo, gere o outline da apresentação.\n\n"
                    f"{fontes}"
                ),
            },
        ],
        temperature=0.3,
        response_format={"type": "json_object"},
    )
    return _extrair_json(resp)


def _aplicar_titulo(shape, texto: str, *, size: int = 28, bold: bool = True) -> None:
    shape.text = texto
    for p in shape.text_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0x00, 0x10, 0x60)
            run.font.name = "Calibri"


def montar_pptx(outline: dict[str, Any], destino: Path) -> Path:
    prs = Presentation()
    # Capa
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.5))
    _aplicar_titulo(box, str(outline.get("titulo") or "Apresentação"), size=36)
    sub = str(outline.get("subtitulo") or "").strip()
    if sub:
        box2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(8.4), Inches(1))
        _aplicar_titulo(box2, sub, size=18, bold=False)

    for item in outline.get("slides") or []:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        titulo = str(item.get("titulo") or "Slide")
        tbox = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.8), Inches(1))
        _aplicar_titulo(tbox, titulo, size=26)
        bullets = item.get("bullets") or []
        body = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
        tf = body.text_frame
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {b}"
            p.level = 0
            for run in p.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(0x00, 0x10, 0x60)
                run.font.name = "Calibri"

    destino.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(destino))
    return destino


def gerar_apresentacao_pptx(caminhos: list[Path]) -> Path:
    ensure_dirs()
    outline = gerar_outline_slides(caminhos)
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    destino = OUTPUTS_DIR / f"apresentacao_{stamp}.pptx"
    return montar_pptx(outline, destino)
