"""Apresentação visual PPTX — skill `apresentacao-visual`.

Deck widescreen 16:9 com estrutura fixa (capa → contexto → ideia →
desenvolvimento → síntese → aplicação → final), paleta canônica e pouco texto.
"""

from __future__ import annotations

import json
import re
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from core.documentos import extrair_texto_arquivo
from core.openai_client import chat_completion, get_api_key
from core.skills_locais import corpo_skill
from core.utils import OUTPUTS_DIR, ensure_dirs

# Widescreen 16:9
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)

FUNDO = RGBColor(0xF7, 0xF8, 0xFC)
TEXTO = RGBColor(0x1F, 0x29, 0x37)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
AZUL = RGBColor(0x3B, 0x82, 0xF6)
VERDE = RGBColor(0x10, 0xB9, 0x81)
AMARELO = RGBColor(0xF5, 0x9E, 0x0B)
ROSA = RGBColor(0xEC, 0x48, 0x99)
ROXO = RGBColor(0x8B, 0x5C, 0xF6)
BORDA = RGBColor(0xE5, 0xE7, 0xEB)

CORES_APOIO = (AZUL, VERDE, AMARELO, ROSA, ROXO)

ICONES = {
    "alvo": "◎",
    "target": "◎",
    "ideia": "✦",
    "lightbulb": "✦",
    "check": "✓",
    "fluxo": "➜",
    "pessoas": "☺",
    "alerta": "!",
    "dado": "▣",
    "estrela": "★",
    "lista": "≡",
    "seta": "→",
    "default": "●",
}

SYSTEM_APRESENTACAO = """\
Você monta o roteiro de uma apresentação executiva visual em {idioma}.
Responda SOMENTE JSON válido neste schema:

{{
  "titulo": "título forte da capa",
  "subtitulo": "uma frase explicativa",
  "contexto": {{
    "titulo": "Contexto",
    "bullets": ["ponto 1", "ponto 2", "ponto 3"]
  }},
  "ideia_central": {{
    "mensagem": "mensagem principal em até 2 frases",
    "destaque": "frase-chave curta"
  }},
  "desenvolvimento": [
    {{
      "titulo": "título curto",
      "icone": "alvo|ideia|check|fluxo|pessoas|alerta|dado|estrela|lista|seta",
      "bullets": ["máx 14 palavras", "..."],
      "destaque": "informação mais importante do slide",
      "layout": "cards|colunas|timeline|fluxo|comparacao"
    }}
  ],
  "sintese": {{
    "bullets": ["aprendizado 1", "..."]
  }},
  "aplicacao": {{
    "bullets": ["como usar...", "..."],
    "proximos_passos": ["passo 1", "passo 2"]
  }},
  "frase_final": "frase-síntese memorável"
}}

REGRAS
- Estrutura obrigatória: capa + contexto + ideia_central + 4 a 8 slides em
  desenvolvimento + sintese + aplicacao + final.
- Máximo 5 bullets por lista; cada bullet ≤ 14 palavras.
- Cada slide = uma ideia. Não invente fatos fora do material.
- Linguagem simples, direta, acionável. Sem "etc.", "robusto", "diversos".
- Prefira números e nomes próprios quando existirem na fonte.
"""


def _set_fill(shape, cor: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor


def _set_line(shape, cor: RGBColor, *, pt: float = 1.0) -> None:
    shape.line.color.rgb = cor
    shape.line.width = Pt(pt)


def _soft_shadow(shape) -> None:
    """Sombra discreta via XML (python-pptx não expõe API simples)."""
    try:
        spPr = shape._element.spPr  # noqa: SLF001
        effect = OxmlElement("a:effectLst")
        outer = OxmlElement("a:outerShdw")
        outer.set("blurRad", "50800")
        outer.set("dist", "38100")
        outer.set("dir", "2700000")
        outer.set("algn", "tl")
        outer.set("rotWithShape", "0")
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", "1F2937")
        alpha = OxmlElement("a:alpha")
        alpha.set("val", "12000")
        srgb.append(alpha)
        outer.append(srgb)
        effect.append(outer)
        spPr.append(effect)
    except Exception:  # noqa: BLE001
        pass


def _textbox(
    slide,
    left,
    top,
    width,
    height,
    texto: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXTO,
    align=PP_ALIGN.LEFT,
    font: str = "Calibri",
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = texto
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _bullets(
    slide,
    left,
    top,
    width,
    height,
    itens: list[str],
    *,
    size: int = 16,
    color: RGBColor = TEXTO,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(itens[:5]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.space_after = Pt(8)
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = "Calibri"


def _card(slide, left, top, width, height, *, fill: RGBColor = BRANCO) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _set_fill(shape, fill)
    _set_line(shape, BORDA, pt=1.0)
    try:
        shape.adjustments[0] = 0.15
    except Exception:  # noqa: BLE001
        pass
    _soft_shadow(shape)
    return shape


def _barra_topo(slide, cor: RGBColor = AZUL) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, _SLIDE_W, Inches(0.12))
    _set_fill(bar, cor)
    bar.line.fill.background()


def _fundo(slide) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, _SLIDE_W, _SLIDE_H)
    _set_fill(bg, FUNDO)
    bg.line.fill.background()
    # Enviar para o fundo
    spTree = slide.shapes._spTree  # noqa: SLF001
    sp = bg._element  # noqa: SLF001
    spTree.remove(sp)
    spTree.insert(2, sp)


def _icone_glyph(nome: str) -> str:
    return ICONES.get((nome or "").strip().lower(), ICONES["default"])


def _extrair_json(texto: str) -> dict[str, Any]:
    bruto = (texto or "").strip()
    cerca = re.search(r"```(?:json)?\s*([\s\S]*?)```", bruto)
    if cerca:
        bruto = cerca.group(1).strip()
    try:
        return json.loads(bruto)
    except json.JSONDecodeError:
        ini, fim = bruto.find("{"), bruto.rfind("}")
        if ini >= 0 and fim > ini:
            return json.loads(bruto[ini : fim + 1])
        raise


def _texto_fontes(caminhos: list[Path], *, max_chars: int = 60000) -> str:
    partes: list[str] = []
    total = 0
    for path in caminhos:
        try:
            doc = extrair_texto_arquivo(path.name, path.read_bytes())
            bloco = f"--- {path.name} ---\n{doc.texto}"
        except Exception as exc:  # noqa: BLE001
            bloco = f"--- {path.name} (erro: {exc}) ---"
        if total + len(bloco) > max_chars:
            resto = max_chars - total
            if resto > 500:
                partes.append(bloco[:resto] + "\n…")
            break
        partes.append(bloco)
        total += len(bloco)
    return "\n\n".join(partes)


def _limpar_bullets(itens: Any) -> list[str]:
    out: list[str] = []
    if not isinstance(itens, list):
        return out
    for raw in itens:
        t = str(raw or "").strip()
        if not t:
            continue
        palavras = t.split()
        if len(palavras) > 14:
            t = " ".join(palavras[:14])
        out.append(t)
        if len(out) >= 5:
            break
    return out


def gerar_roteiro_apresentacao(
    caminhos: list[Path],
    especificacoes: str = "",
    *,
    idioma: str = "pt-BR",
) -> dict[str, Any]:
    if not get_api_key():
        raise RuntimeError("Chave de API do provedor LLM não configurada.")
    from core.especificacoes_llm import anexar_especificacoes

    fontes = _texto_fontes(caminhos)
    if not fontes.strip():
        raise ValueError("Fontes vazias.")
    user = anexar_especificacoes(
        "Antes de montar o JSON:\n"
        "1. Analise todo o conteúdo.\n"
        "2. Extraia os conceitos principais.\n"
        "3. Organize a narrativa da estrutura fixa.\n"
        "4. Defina 4 a 8 slides de desenvolvimento.\n"
        "5. Só então devolva o JSON.\n\n"
        f"### Material\n\n{fontes}",
        especificacoes,
    )
    skill_body = corpo_skill("apresentacao-visual")
    system = SYSTEM_APRESENTACAO.format(idioma=idioma)
    if skill_body:
        system = f"{system}\n\n### Skill `apresentacao-visual`\n{skill_body}"
    raw = chat_completion(
        [
            {"role": "system", "content": system},
            {"role": "user", "content": user[:120000]},
        ],
        temperature=0.3,
        response_format={"type": "json_object"},
    )
    return _extrair_json(raw)


def _slide_capa(prs: Presentation, dados: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fundo(slide)
    _barra_topo(slide, AZUL)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(0.35), Inches(2.2)
    )
    _set_fill(accent, AZUL)
    accent.line.fill.background()
    _textbox(
        slide,
        Inches(1.5),
        Inches(2.1),
        Inches(10.5),
        Inches(1.4),
        str(dados.get("titulo") or "Apresentação"),
        size=40,
        bold=True,
    )
    sub = str(dados.get("subtitulo") or "").strip()
    if sub:
        _textbox(
            slide,
            Inches(1.5),
            Inches(3.6),
            Inches(10.5),
            Inches(1.0),
            sub,
            size=20,
            color=RGBColor(0x4B, 0x55, 0x63),
        )
    badge = _card(slide, Inches(1.5), Inches(5.2), Inches(4.2), Inches(0.7), fill=AZUL)
    _textbox(
        slide,
        Inches(1.7),
        Inches(5.3),
        Inches(3.8),
        Inches(0.5),
        "Resumo executivo visual",
        size=14,
        bold=True,
        color=BRANCO,
        align=PP_ALIGN.CENTER,
    )
    _ = badge


def _slide_contexto(prs: Presentation, bloco: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fundo(slide)
    _barra_topo(slide, VERDE)
    _textbox(
        slide,
        Inches(0.8),
        Inches(0.45),
        Inches(11.5),
        Inches(0.7),
        str(bloco.get("titulo") or "Contexto"),
        size=28,
        bold=True,
    )
    bullets = _limpar_bullets(bloco.get("bullets"))
    n = max(1, len(bullets))
    largura = Inches(11.5) / n
    for i, b in enumerate(bullets):
        left = Inches(0.8) + largura * i
        cor = CORES_APOIO[i % len(CORES_APOIO)]
        _card(slide, left, Inches(1.5), largura - Inches(0.2), Inches(4.8))
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, Inches(1.5), largura - Inches(0.2), Inches(0.12)
        )
        _set_fill(bar, cor)
        bar.line.fill.background()
        _textbox(
            slide,
            left + Inches(0.2),
            Inches(1.9),
            largura - Inches(0.5),
            Inches(4.0),
            b,
            size=15,
            bold=True,
        )


def _slide_ideia(prs: Presentation, bloco: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fundo(slide)
    _barra_topo(slide, ROXO)
    _textbox(
        slide,
        Inches(0.8),
        Inches(0.45),
        Inches(11.5),
        Inches(0.6),
        "Ideia central",
        size=28,
        bold=True,
    )
    _card(slide, Inches(1.2), Inches(1.6), Inches(10.9), Inches(3.2), fill=BRANCO)
    msg = str(bloco.get("mensagem") or "").strip()
    _textbox(
        slide,
        Inches(1.6),
        Inches(2.0),
        Inches(10.1),
        Inches(2.2),
        msg,
        size=24,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    dest = str(bloco.get("destaque") or "").strip()
    if dest:
        _card(slide, Inches(2.5), Inches(5.2), Inches(8.3), Inches(1.0), fill=ROXO)
        _textbox(
            slide,
            Inches(2.8),
            Inches(5.35),
            Inches(7.7),
            Inches(0.7),
            dest,
            size=16,
            bold=True,
            color=BRANCO,
            align=PP_ALIGN.CENTER,
        )


def _slide_desenvolvimento(prs: Presentation, item: dict[str, Any], idx: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fundo(slide)
    cor = CORES_APOIO[idx % len(CORES_APOIO)]
    _barra_topo(slide, cor)
    glyph = _icone_glyph(str(item.get("icone") or "default"))
    disc = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), Inches(0.4), Inches(0.65), Inches(0.65)
    )
    _set_fill(disc, cor)
    disc.line.fill.background()
    _textbox(
        slide,
        Inches(0.8),
        Inches(0.48),
        Inches(0.65),
        Inches(0.5),
        glyph,
        size=18,
        bold=True,
        color=BRANCO,
        align=PP_ALIGN.CENTER,
    )
    _textbox(
        slide,
        Inches(1.7),
        Inches(0.45),
        Inches(10.5),
        Inches(0.7),
        str(item.get("titulo") or f"Tópico {idx + 1}"),
        size=26,
        bold=True,
    )
    bullets = _limpar_bullets(item.get("bullets"))
    _card(slide, Inches(0.8), Inches(1.4), Inches(8.0), Inches(4.5))
    _bullets(slide, Inches(1.1), Inches(1.7), Inches(7.4), Inches(3.9), bullets, size=17)
    dest = str(item.get("destaque") or "").strip()
    if dest:
        _card(slide, Inches(9.1), Inches(1.4), Inches(3.4), Inches(4.5), fill=cor)
        _textbox(
            slide,
            Inches(9.35),
            Inches(1.7),
            Inches(2.9),
            Inches(0.5),
            "Destaque",
            size=12,
            bold=True,
            color=BRANCO,
        )
        _textbox(
            slide,
            Inches(9.35),
            Inches(2.4),
            Inches(2.9),
            Inches(3.0),
            dest,
            size=16,
            bold=True,
            color=BRANCO,
        )


def _slide_lista(
    prs: Presentation,
    titulo: str,
    bullets: list[str],
    *,
    cor: RGBColor,
    extras: list[str] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fundo(slide)
    _barra_topo(slide, cor)
    _textbox(slide, Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.7), titulo, size=28, bold=True)
    _card(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.2))
    _bullets(slide, Inches(1.2), Inches(1.8), Inches(11.0), Inches(3.2), bullets, size=18)
    if extras:
        y = Inches(5.0)
        _textbox(
            slide,
            Inches(1.2),
            y,
            Inches(11.0),
            Inches(0.4),
            "Próximos passos",
            size=14,
            bold=True,
            color=cor,
        )
        _bullets(slide, Inches(1.2), Inches(5.4), Inches(11.0), Inches(1.0), extras, size=15)


def _slide_final(prs: Presentation, frase: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fundo(slide)
    _barra_topo(slide, ROSA)
    _card(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(2.8), fill=BRANCO)
    _textbox(
        slide,
        Inches(2.0),
        Inches(2.7),
        Inches(9.3),
        Inches(1.8),
        frase or "Obrigado.",
        size=26,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _textbox(
        slide,
        Inches(2.0),
        Inches(5.5),
        Inches(9.3),
        Inches(0.5),
        "Apresentação visual · pronta para uso",
        size=12,
        color=RGBColor(0x6B, 0x72, 0x80),
        align=PP_ALIGN.CENTER,
    )


def montar_pptx_visual(dados: dict[str, Any], destino: Path) -> Path:
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H

    _slide_capa(prs, dados)
    _slide_contexto(prs, dados.get("contexto") or {})
    _slide_ideia(prs, dados.get("ideia_central") or {})

    desenvolvimento = dados.get("desenvolvimento") or []
    if not isinstance(desenvolvimento, list):
        desenvolvimento = []
    for i, item in enumerate(desenvolvimento[:8]):
        if isinstance(item, dict):
            _slide_desenvolvimento(prs, item, i)

    sintese = dados.get("sintese") or {}
    _slide_lista(
        prs,
        "Síntese",
        _limpar_bullets(sintese.get("bullets") if isinstance(sintese, dict) else []),
        cor=VERDE,
    )

    aplicacao = dados.get("aplicacao") or {}
    if not isinstance(aplicacao, dict):
        aplicacao = {}
    _slide_lista(
        prs,
        "Aplicação prática",
        _limpar_bullets(aplicacao.get("bullets")),
        cor=AMARELO,
        extras=_limpar_bullets(aplicacao.get("proximos_passos")),
    )

    _slide_final(prs, str(dados.get("frase_final") or "").strip())

    destino.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(destino))
    return destino


def gerar_apresentacao_visual(
    caminhos: list[Path],
    especificacoes: str = "",
    *,
    idioma: str = "pt-BR",
    nome_arquivo: str | None = None,
) -> Path:
    """LLM (roteiro estruturado) → PPTX 16:9 com a identidade da skill."""
    ensure_dirs()
    if not caminhos:
        raise ValueError("Nenhuma fonte selecionada.")
    dados = gerar_roteiro_apresentacao(
        caminhos, especificacoes=especificacoes, idioma=idioma
    )
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    nome = nome_arquivo or f"apresentacao_visual_{stamp}.pptx"
    destino = OUTPUTS_DIR / nome
    return montar_pptx_visual(dados, destino)
